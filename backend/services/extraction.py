import fitz  # pymupdf
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


def extract_text(file_path: str, mime_type: str) -> list[dict]:
    """Extract text from a document. Returns list of {text, page_number}."""
    if mime_type == "application/pdf":
        return _extract_pdf(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {mime_type}")


def _extract_pdf(file_path: str) -> list[dict]:
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            # We track character range within the page. 
            # In a full RAG we might track global offsets across the document.
            pages.append({
                "text": text, 
                "page_number": i + 1,
                "start_offset": 0,
                "end_offset": len(text)
            })
    return pages


def _extract_docx(file_path: str) -> list[dict]:
    doc = DocxDocument(file_path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": text, "page_number": None}]


def _extract_xlsx(file_path: str) -> list[dict]:
    wb = load_workbook(file_path, read_only=True, data_only=True)
    pages = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            pages.append({"text": f"Sheet: {sheet.title}\n" + "\n".join(rows), "page_number": None})
    return pages


def _extract_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            pages.append({"text": "\n".join(texts), "page_number": i + 1})
    return pages


def get_page_count(file_path: str, mime_type: str) -> int | None:
    if mime_type == "application/pdf":
        doc = fitz.open(file_path)
        return len(doc)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(file_path)
        return len(prs.slides)
    return None
